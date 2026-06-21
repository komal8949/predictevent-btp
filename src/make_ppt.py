"""Generate the finale slide deck -> outputs/PredictEvent_Deck.pptx"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "outputs")

NAVY = RGBColor(0x1A, 0x36, 0x5D)
BLUE = RGBColor(0x2B, 0x6C, 0xB0)
GREY = RGBColor(0x44, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def box(slide, l, t, w, h, text, size=18, bold=False, color=GREY,
        align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb


def band(slide, color=NAVY, h=Inches(1.1)):
    s = slide.shapes.add_shape(1, 0, 0, W, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def header(slide, title):
    band(slide)
    box(slide, Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.7),
        title, size=28, bold=True, color=WHITE)


def bullets(slide, items, l=Inches(0.7), t=Inches(1.4), w=Inches(12), size=18):
    tb = slide.shapes.add_textbox(l, t, w, Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - lvl * 2)
        r.font.color.rgb = NAVY if lvl == 0 else GREY
        r.font.bold = (lvl == 0)
        p.space_after = Pt(6)


def img(slide, path, l, t, w):
    if os.path.exists(os.path.join(OUT, path)):
        slide.shapes.add_picture(os.path.join(OUT, path), l, t, width=w)


# ---- Slide 1: title ----
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
box(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.2),
    "PredictEvent", size=54, bold=True, color=WHITE)
box(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.8),
    "Event-Driven Congestion Intelligence for Bengaluru Traffic Police",
    size=24, color=RGBColor(0xCB, 0xD5, 0xE0))
box(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.6),
    "Theme 2  |  BTP × Flipkart Hackathon  |  Forecast impact → recommend manpower, barricading & diversion",
    size=15, color=RGBColor(0x9F, 0xB3, 0xC8))

# ---- Slide 2: problem ----
s = prs.slides.add_slide(BLANK); header(s, "The Problem")
bullets(s, [
    ("Rallies, festivals, sports, processions, VIP movement, construction & sudden incidents choke Bengaluru's corridors.", 0),
    ("Event impact is not quantified in advance — which corridor, how hard, tomorrow?", 0),
    ("Resource deployment is experience-driven — manpower & barricades decided by gut feel.", 0),
    ("No post-event learning loop — the same mistakes repeat.", 0),
    ("Goal: turn 8,173 historical events into a forward-looking deployment planner.", 0),
])

# ---- Slide 3: solution overview ----
s = prs.slides.add_slide(BLANK); header(s, "Our Solution")
bullets(s, [
    ("A laptop-deployable decision-support engine (no GPU) with 3 parts:", 0),
    ("Hotspot Load Forecaster  →  forecast load per corridor for tomorrow", 1),
    ("Closure-Risk Classifier  →  barricade & RED/AMBER/GREEN triage", 1),
    ("OR Optimization Engine  →  SLA-guaranteed officers + coverage-optimal placement", 1),
    ("Uses ONLY the provided ASTRAM dataset. Interpretable, deployable today.", 0),
], size=20)
img(s, "spatial_scatter.png", Inches(8.6), Inches(1.5), Inches(4.4))

# ---- Slide 4: ML results ----
s = prs.slides.add_slide(BLANK); header(s, "ML Results — honest, time-based split")
rows = [
    ("Component", "Result", "Baseline"),
    ("Hotspot Load Forecaster", "R² 0.536 · MAE 2.19", "seasonal-naive R² −0.32"),
    ("Closure-Risk Classifier", "ROC-AUC 0.813 · recall 0.76", "random PR-AUC 0.072"),
    ("Duration Predictor", "not reliably learnable", "(data limitation, see findings)"),
]
tbl = s.shapes.add_table(4, 3, Inches(0.6), Inches(1.5), Inches(12.1), Inches(2.6)).table
for c in range(3):
    tbl.columns[c].width = Inches([3.8, 4.4, 3.9][c])
for r in range(4):
    for c in range(3):
        cell = tbl.cell(r, c); cell.text = rows[r][c]
        p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(15)
        if r == 0:
            p.runs[0].font.bold = True; p.runs[0].font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
        else:
            p.runs[0].font.color.rgb = GREY
box(s, Inches(0.6), Inches(4.35), Inches(8.0), Inches(2.6),
    "Operating point: recall-floor — catches ≥75% of closures (a missed closure costs BTP more).\n"
    "Top drivers (permutation importance): event_cause ≫ corridor > location > vehicle type.\n"
    "Deep ST-model: MAE 2.12 vs GBM 2.24 (~5% better, R² tied); ensemble R² 0.55. GBM kept in prod.\n"
    "Calibrated probabilities: ECE 0.33→0.02, Brier 0.185→0.058 — '81% risk' really means ~81%.\n"
    "Cross-validated (5-fold expanding): hotspot R² 0.56±0.14, closure AUC 0.73±0.05 — spread, not cherry-pick.",
    size=14, color=NAVY)
img(s, "closure_calibration.png", Inches(8.9), Inches(4.3), Inches(4.0))

# ---- Slide 4b: OR Optimization Engine (the differentiator) ----
s = prs.slides.add_slide(BLANK); header(s, "OR Optimization Engine — the differentiator")
bullets(s, [
    ("We replaced heuristic rules with the forecast→optimize pattern used in deployed police systems.", 0),
    ("Manpower = M/M/c queueing (Erlang-C): officer count sized so expected response wait ≤ SLA (e.g. 5 min) — a number WITH a service-level guarantee.", 0),
    ("Barricade placement = Maximum Coverage Location Problem (MCLP): place P units to cover max risk×load-weighted demand (PuLP/CBC).", 0),
    ("Risk×load weighting: both high-traffic AND high-risk corridors, per OR location-allocation literature.", 0),
    ("Every number is traceable — no black box. Evidence: SMU IJCAI-2019, Vlahogianni 2019, Miao/Easa 2021.", 0),
], size=17)

# ---- Slide 5: data findings (the trust slide) ----
s = prs.slides.add_slide(BLANK); header(s, "Data Findings That Build Trust")
bullets(s, [
    ("Timestamps are local IST mislabeled '+00' — proven via bimodal rush-hour signature; we use stored hour as-is.", 0),
    ("closed_datetime is administrative ticket-close, NOT on-ground clearance → exact duration not learnable. We refuse to over-claim it.", 0),
    ("We caught & removed target leakage: end-coordinate flag is filled BECAUSE a closure happened (98% vs 0%) — faked ROC-AUC 1.0.", 0),
    ("Permutation importance exposed it; dropping it gives the honest 0.81.", 1),
    ("Event-driven causes (public_event, procession, VIP, protest) have 40–80% closure rates — the true Theme-2 targets.", 0),
], size=17)

# ---- Slide 6: prototype ----
s = prs.slides.add_slide(BLANK); header(s, "Prototype — Live Demo")
bullets(s, [
    ("Event Simulator — type an event (e.g. IPL match @ Chinnaswamy) → instant risk %, tier, personnel, barricade & diversion.", 0),
    ("Hotspot Forecast + Manpower — pick a day → corridor load + recommended personnel.", 0),
    ("Live Event Triage — every event ranked by closure-risk with action tier + map.", 0),
    ("Analytics — historical patterns + feature importance.", 0),
], size=18, t=Inches(1.4))
img(s, "events_by_hour.png", Inches(0.7), Inches(4.0), Inches(6.0))
img(s, "daily_timeline.png", Inches(6.9), Inches(4.0), Inches(6.0))

# ---- Slide 6b: BTP stack fit (verified positioning) ----
s = prs.slides.add_slide(BLANK); header(s, "Where It Fits in BTP's Stack")
bullets(s, [
    ("BATCS optimizes the SIGNALS — adaptive timing at 165 junctions (CoSiCoSt). Not people.", 0),
    ("ASTRAM SEES the city — awareness, 15-min alerts, microsimulation. (Our dataset's source.)", 0),
    ("Neither produces a prescriptive police-deployment plan.", 0),
    ("PredictEvent DECIDES the deployment — turns ASTRAM's event data into an officer count with a service-level guarantee + coverage-optimal placement.", 0),
    ("A prescriptive module on top of the stack BTP already runs — it extends ASTRAM, doesn't compete.", 0),
], size=18)

# ---- Slide 7: why it wins + roadmap ----
s = prs.slides.add_slide(BLANK); header(s, "Why It Wins · Roadmap")
bullets(s, [
    ("Solves BTP's actual #1 pain — event manpower planning, not a tech demo.", 0),
    ("OR-optimized, SLA-guaranteed deployment — not heuristic rules.", 0),
    ("Honest, leakage-checked metrics — survives a technical jury.", 0),
    ("Plugs into ASTRAM; deployable today; provided-data only.", 0),
    ("Roadmap:", 0),
    ("Live ASTRAM API feed → real clearance timestamps", 1),
    ("Capacitated city-wide MCLP + IP shift-scheduling (ILS/tabu fast replanning)", 1),
    ("Mobile officer app + automatic retraining loop", 1),
])

prs.save(os.path.join(OUT, "PredictEvent_Deck.pptx"))
print("Saved", os.path.join(OUT, "PredictEvent_Deck.pptx"), "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
